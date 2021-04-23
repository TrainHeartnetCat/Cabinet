# coding: utf-8
from pptx import Presentation
import copy
import datetime
import os

print('please run in python3')

# 需要合并的pptx所在文件夹
PATH_PPTX = r"C:\Users\OldKuroCat\Desktop\pptcombine\test\target"
# 合并后得到的pptx输出文件夹
PATH_TEMP = r"C:\Users\OldKuroCat\Desktop\pptcombine\test\output"
# 合并模板pptx，注意不能为空pptx
PATH_BASE = r'C:\Users\OldKuroCat\Desktop\pptcombine\test\base.pptx'


# 变量分别是base pptx路径，需要合并的pptx路径，需要合并的pptx中具体slide编号，
def copy_slide(prs_base, prs_add, add_index_no, base_master_no, base_layout_no):
    source = prs_add.slides[int(add_index_no)]
    dest = prs_base.slides.add_slide(prs_base.slide_masters[int(base_master_no)].slide_layouts[int(base_layout_no)])

    for placeholder in dest.shapes.placeholders:
        if placeholder.has_text_frame and placeholder.text_frame.text == "":
            sp = placeholder._sp
            sp.getparent().remove(sp)

    for shape in source.shapes:
        newel = copy.deepcopy(shape.element)
        dest.shapes._spTree.insert_element_before(newel, 'p:extLst')

    for key, value in source.part.rels.items():
        if not "notesSlide" in value.reltype and not "slideLayout" in value.reltype:
            dest.part.rels.add_relationship(value.reltype, value._target, value.rId)
    return dest



prs_base = Presentation(PATH_BASE)

pptx_list = os.listdir(PATH_PPTX)

for ppt in pptx_list:
    prs_cu = Presentation(PATH_PPTX + '/' + ppt)
    copy_slide(prs_base, prs_cu, 0, 0, 0)
    del prs_cu

save_name = "merged" + '_{0:%Y%m%d%H%M%S}'.format(datetime.datetime.now()) + ".pptx"
prs_base.save(PATH_TEMP + '/' + save_name)

print("work complete!")
print("The merged pptx path is")
print(PATH_TEMP + '/' + save_name)



